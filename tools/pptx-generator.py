#!/usr/bin/env python3
"""
tools/pptx-generator.py — Insights Forge deck generator

Accepts a JSON spec and produces a branded .pptx using
Dynatrace_Brand_Insights-Forge.pptx as the template base.
All brand mechanics (layouts, fonts, master, footer) are handled
by the template. This tool fills in engagement-specific content.

By default the generator performs a read-only check of the DT Flow fonts:
it compares DTFlow/ against the user font directory and prints a one-line
notice if any are missing, but never installs anything on its own.
Font installation is strictly opt-in via the --install-fonts flag.

Usage:
    python3 tools/pptx-generator.py <spec.json> [output.pptx]
        Generate a deck. Output defaults to <spec-dir>/deck-YYYY-MM-DD.pptx

    python3 tools/pptx-generator.py --list-layouts
        Print all 64 named layouts in the template.

    python3 tools/pptx-generator.py --install-fonts
        Install DT Flow fonts from DTFlow/ to the user font directory
        (macOS: ~/Library/Fonts, Linux: ~/.fonts). This is the only mode
        that writes to the font directory. Safe to run multiple times —
        skips fonts already installed. Windows has no auto-install path:
        right-click each DTFlow/*.otf and choose Install.

Spec format:  tools/pptx-spec-example.json

Requires Python 3.9+ and python-pptx (pip install -r tools/requirements.txt).
"""

import sys
import re
import json
import platform
import shutil
from datetime import date
from pathlib import Path
from typing import Optional
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

PROJECT_ROOT = Path(__file__).resolve().parent.parent
TEMPLATE     = PROJECT_ROOT / "Dynatrace_Brand_Insights-Forge.pptx"
FONTS_DIR    = PROJECT_ROOT / "DTFlow"
ASSETS_DIR   = PROJECT_ROOT / "assets"

# Brand constants — brand-spec.md §2 (theme colors) and §5 (chart series)
_OVERLAY_COLOR = RGBColor(0x1A, 0x24, 0x40)   # Deep navy #1A2440
_WHITE         = RGBColor(0xFF, 0xFF, 0xFF)   # Light 1 — text over the overlay
_FOOTER_GRAY   = RGBColor(0x6F, 0x74, 0x7F)   # Light 2 — brand-spec §8 footer
# Layouts that carry no footer in the template (cover / closing) — a wave
# slide on one of these does not get the re-added footer line either.
_NO_FOOTER_LAYOUT_PREFIXES = ("Title Slide", "Title slide", "Thank you")

# Series order follows brand-spec.md §5 (the chart-series canon), which is
# deliberately NOT numeric accent order: 1→A1, 2→A3, 3→A5, 4→A4, 5→A6, 6→A2.
BRAND_CHART_COLORS = [
    RGBColor(0x49, 0xC2, 0xB3),  # Series 1 — Teal       (Accent 1)
    RGBColor(0x19, 0x66, 0xFF),  # Series 2 — Royal blue (Accent 3)
    RGBColor(0x8D, 0x1C, 0xDC),  # Series 3 — Violet     (Accent 5)
    RGBColor(0x5E, 0x28, 0xE5),  # Series 4 — Purple     (Accent 4)
    RGBColor(0xC9, 0x3F, 0xDB),  # Series 5 — Magenta    (Accent 6)
    RGBColor(0x3B, 0xAC, 0xF0),  # Series 6 — Light blue (Accent 2)
]

# Named wave asset shortcuts (files live in assets/)
WAVE_ASSETS = {
    "wave-bg":  ASSETS_DIR / "wave-bg.png",   # cover / closing
    "wave-ask": ASSETS_DIR / "wave-ask.png",  # decision-required accent slides
}


# ── Font installation ──────────────────────────────────────────────────────

def _user_font_dir() -> Optional[Path]:
    """Return the user font directory for the current OS, or None if unsupported.

    Windows is deliberately unsupported: per-user fonts there need a registry
    entry as well as the file copy, and half an install is worse than none.
    """
    system = platform.system()
    if system == "Darwin":
        return Path.home() / "Library" / "Fonts"
    if system == "Linux":
        return Path.home() / ".fonts"
    return None


def ensure_dtflow_fonts(verbose: bool = True) -> int:
    """Install DT Flow fonts from DTFlow/ into the user font directory.

    Copies only fonts that are not already present.  Safe to call on
    every generator run — the check is fast and the copy only happens
    once.  Returns the number of fonts newly installed.
    """
    if not FONTS_DIR.exists():
        if verbose:
            print(f"  ⚠ DTFlow/ not found — skipping font install", file=sys.stderr)
        return 0

    font_dir = _user_font_dir()
    if font_dir is None:
        if verbose:
            print(f"  ⚠ Font auto-install not supported on {platform.system()}."
                  f"  Install DTFlow/*.otf manually.", file=sys.stderr)
        return 0

    font_dir.mkdir(parents=True, exist_ok=True)

    otf_files = sorted(FONTS_DIR.glob("*.otf"))
    if not otf_files:
        if verbose:
            print("  ⚠ No .otf files found in DTFlow/", file=sys.stderr)
        return 0

    installed = []
    for otf in otf_files:
        dest = font_dir / otf.name
        if not dest.exists():
            shutil.copy2(otf, dest)
            installed.append(otf.name)

    if installed:
        print(f"  ✓ Installed {len(installed)} DT Flow font(s) → {font_dir}")
        # Refresh font cache on Linux so apps can see the new fonts
        if platform.system() == "Linux":
            import subprocess
            subprocess.run(["fc-cache", "-fv"], capture_output=True)
    elif verbose:
        print(f"  ✓ DT Flow fonts already installed ({len(otf_files)} fonts in {font_dir})")

    return len(installed)


def check_dtflow_fonts() -> None:
    """Read-only check of the DT Flow fonts — performs ZERO writes.

    Compares the .otf filenames in DTFlow/ against the user font directory.
    If the font directory does not exist (it is never created here), every
    font counts as missing. Silent when all fonts are present; otherwise
    prints exactly one notice pointing at --install-fonts. Installation
    itself happens only via ensure_dtflow_fonts() behind that flag.
    """
    if not FONTS_DIR.exists():
        return
    otf_files = sorted(FONTS_DIR.glob("*.otf"))
    if not otf_files:
        return

    font_dir = _user_font_dir()
    if font_dir is None or not font_dir.is_dir():
        missing = len(otf_files)
    else:
        present = {p.name for p in font_dir.iterdir()}
        missing = sum(1 for otf in otf_files if otf.name not in present)

    if missing:
        how = ("Run with --install-fonts to install them."
               if font_dir is not None
               else f"Install DTFlow/*.otf manually on {platform.system()}.")
        print(f"⚠ {missing} DT Flow font(s) not installed — deck will render "
              f"with fallback fonts. {how}")


# ── Visual helpers — wave backgrounds, overlays, chart colors ─────────────

def _slide_dimensions(slide):
    """Return (width, height) in EMUs from the presentation."""
    prs = slide.part.package.presentation_part.presentation
    return prs.slide_width, prs.slide_height


def _move_shape_in_spTree(slide, shape, position: int) -> None:
    """Reposition a shape element in the slide's spTree to control z-order.

    spTree children 0–1 are <p:nvGrpSpPr> and <p:grpSpPr>, so position=2
    is the bottom of the shape stack (renders behind every other shape).
    """
    spTree = slide.shapes._spTree
    spTree.remove(shape._element)
    spTree.insert(position, shape._element)


def _set_fill_alpha(shape, opacity: float) -> None:
    """Set opacity on a shape whose fill is already set to solid().

    opacity: 0.0 (transparent) → 1.0 (fully opaque).
    OOXML alpha is expressed in thousandths-of-a-percent: 100% = 100000.
    """
    alpha_val = str(int(opacity * 100000))
    spPr = shape._element.spPr
    solidFill = spPr.find(".//" + qn("a:solidFill"))
    if solidFill is None:
        return
    # Try srgbClr first, then sysClr / schemeClr
    for tag in (qn("a:srgbClr"), qn("a:sysClr"), qn("a:schemeClr")):
        clr = solidFill.find(tag)
        if clr is not None:
            # Remove any existing alpha child, then add a fresh one
            for existing in clr.findall(qn("a:alpha")):
                clr.remove(existing)
            alpha_elem = etree.SubElement(clr, qn("a:alpha"))
            alpha_elem.set("val", alpha_val)
            return


def add_wave_background(slide, wave: str, overlay_opacity: float = 0.80) -> None:
    """Insert a wave PNG as a full-slide background with a dark overlay.

    wave: a key from WAVE_ASSETS ("wave-bg", "wave-ask") or a path string.
    overlay_opacity: 0.0–1.0. Default 0.80 (body-text dark slides).
                     Use 0.65–0.70 for title-only dark slides (cover, closing).
                     Use 0.80–0.85 for dark slides with body text.

    Z-order after this call (bottom → top):
        [0] wave PNG picture
        [1] dark overlay rectangle
        [2+] all existing slide content shapes (placeholders etc.)
        [top] re-added footer line (body layouts only)

    The layout's own text on a body layout is black, which vanishes on the
    navy overlay, so every text run on the slide is set to white; and the
    layout/master footer sits BELOW the overlay in z-order, so the
    brand-spec §8 footer line is re-added on top for body layouts.
    """
    # Resolve the wave asset path first so a missing file skips cleanly
    if wave in WAVE_ASSETS:
        wave_path = WAVE_ASSETS[wave]
    else:
        wave_path = Path(wave)
        if not wave_path.is_absolute():
            wave_path = PROJECT_ROOT / wave_path

    if not wave_path.exists():
        print(f"  ⚠ wave asset not found: {wave_path} — skipping wave background",
              file=sys.stderr)
        return

    slide_w, slide_h = _slide_dimensions(slide)
    overlay_opacity = min(max(float(overlay_opacity), 0.0), 1.0)

    # 1. Wave PNG — inserted last so we can move it immediately
    pic = slide.shapes.add_picture(str(wave_path), 0, 0, slide_w, slide_h)
    _move_shape_in_spTree(slide, pic, 2)  # bottom of the shape stack

    # 2. Dark overlay — full-slide rectangle with semi-transparency
    # Use a textbox as the rectangle primitive (cleanest in python-pptx)
    overlay = slide.shapes.add_textbox(Emu(0), Emu(0), slide_w, slide_h)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = _OVERLAY_COLOR
    overlay.line.fill.background()  # no border
    _set_fill_alpha(overlay, overlay_opacity)
    _move_shape_in_spTree(slide, overlay, 3)  # just above the wave image

    # 3. Legibility — white text over the overlay, footer re-added on top
    _restyle_for_dark_background(slide, exclude=(pic, overlay))
    if not slide.slide_layout.name.startswith(_NO_FOOTER_LAYOUT_PREFIXES):
        _add_dark_slide_footer(slide, slide_w, slide_h)


def _restyle_for_dark_background(slide, exclude=()) -> None:
    """Set every text run on the slide to white so body-layout text (black
    by default) stays legible over the dark overlay."""
    skip = {id(s._element) for s in exclude}
    for shape in slide.shapes:
        if id(shape._element) in skip or not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = _WHITE


def _add_dark_slide_footer(slide, slide_w, slide_h) -> None:
    """Re-add the brand-spec §8 footer line above the overlay (the layout's
    footer is covered). Light 2 gray, lower-left, current year."""
    box = slide.shapes.add_textbox(Inches(0.45), slide_h - Inches(0.5),
                                   Inches(6.0), Inches(0.3))
    tf = box.text_frame
    tf.text = f"© {date.today().year} Dynatrace, LLC.   Confidential"
    for run in tf.paragraphs[0].runs:
        run.font.size = Pt(9)
        run.font.color.rgb = _FOOTER_GRAY


def apply_brand_chart_colors(chart) -> None:
    """Apply Dynatrace brand series colors (Accent 1–6) to every series.

    Removes the Office theme color reference so the explicit RGB wins
    regardless of what theme is embedded in the template.

    Line-family charts (LINE*, XY_SCATTER_LINES*, RADAR*) draw their series
    color from the line element (a:ln), not the shape fill, so those get
    the brand color on the line as well.
    """
    is_line_family = any(tag in str(chart.chart_type)
                         for tag in ("LINE", "RADAR"))
    for i, series in enumerate(chart.series):
        color = BRAND_CHART_COLORS[i % len(BRAND_CHART_COLORS)]
        try:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color
            if is_line_family:
                series.format.line.color.rgb = color

            # Clear the theme color reference (<c:spPr><a:solidFill><a:schemeClr>)
            # so the explicit srgbClr we just wrote is not overridden by the theme.
            spPr = series._element.find(qn("c:spPr"))
            if spPr is not None:
                for sf in spPr.findall(".//" + qn("a:solidFill")):
                    for sc in sf.findall(qn("a:schemeClr")):
                        sf.remove(sc)
        except Exception as e:
            print(f"  ⚠ apply_brand_chart_colors series {i}: {e}", file=sys.stderr)


# ── Template helpers ───────────────────────────────────────────────────────

def load_template() -> Presentation:
    if not TEMPLATE.exists():
        raise FileNotFoundError(f"Template not found: {TEMPLATE}")
    return Presentation(str(TEMPLATE))


def list_layouts(prs: Presentation) -> list[str]:
    return [l.name for l in prs.slide_masters[0].slide_layouts]


def get_layout(prs: Presentation, name: str):
    for layout in prs.slide_masters[0].slide_layouts:
        if layout.name == name:
            return layout
    available = "\n  ".join(list_layouts(prs))
    raise ValueError(f"Layout '{name}' not found.\nAvailable:\n  {available}")


def clear_sample_slides(prs: Presentation) -> None:
    """Remove all sample slides from the template, keeping master + layouts."""
    sldIdLst = prs.slides._sldIdLst
    rId_attr = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    rIds = [el.get(rId_attr) for el in list(sldIdLst)]
    for rId in rIds:
        if rId:
            try:
                prs.part.drop_rel(rId)
            except KeyError:
                pass
    sldIdLst.clear()


def add_slide(prs: Presentation, layout_name: str):
    """Add a slide and return (slide, name_to_idx).

    name_to_idx maps each layout placeholder's semantic name (lowercase)
    to its stable placeholder format index. Filling by idx is reliable;
    filling by instance name is not (instance names differ from layout names).
    """
    layout = get_layout(prs, layout_name)
    slide = prs.slides.add_slide(layout)
    name_to_idx = {
        ph.name.lower().strip(): ph.placeholder_format.idx
        for ph in layout.placeholders
    }
    return slide, name_to_idx


# ── Placeholder fill helpers ───────────────────────────────────────────────

def _get_ph(slide, ph_name: str, name_to_idx: dict):
    """Return the placeholder matching ph_name, or None."""
    # Primary: look up idx from the layout's semantic name, then find by idx
    idx = name_to_idx.get(ph_name.lower().strip())
    if idx is not None:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == idx:
                return ph
    # Fallback: direct instance-name match
    for ph in slide.placeholders:
        if ph.name.lower().strip() == ph_name.lower().strip():
            return ph
    return None


def fill(slide, ph_name: str, text: str, n2i: dict) -> bool:
    """Set plain text on a named placeholder. Returns True if found."""
    ph = _get_ph(slide, ph_name, n2i)
    if ph is None:
        return False
    try:
        ph.text = str(text)
    except Exception as e:
        print(f"    ⚠ fill('{ph_name}'): {e}", file=sys.stderr)
    return True


def fill_bullets(slide, ph_name: str, items: list, n2i: dict) -> bool:
    """Fill a placeholder with a bulleted list. Returns True if found."""
    ph = _get_ph(slide, ph_name, n2i)
    if ph is None:
        return False
    if not items:
        return True
    try:
        tf = ph.text_frame
        tf.text = str(items[0])
        for item in items[1:]:
            p = tf.add_paragraph()
            p.text = str(item)
    except Exception as e:
        print(f"    ⚠ fill_bullets('{ph_name}'): {e}", file=sys.stderr)
    return True


def fill_first(slide, candidates: list[str], value, n2i: dict) -> bool:
    """Try placeholder names in order; fill the first one found."""
    for name in candidates:
        if _get_ph(slide, name, n2i) is not None:
            if isinstance(value, list):
                return fill_bullets(slide, name, value, n2i)
            return fill(slide, name, value, n2i)
    return False


def fill_content(slide, ph_name: str, value, n2i: dict) -> bool:
    """Fill a content placeholder — handles both str and list[str]."""
    if isinstance(value, list):
        return fill_bullets(slide, ph_name, value, n2i)
    return fill(slide, ph_name, str(value), n2i)


def _blank_unused_card_slots(slide, n2i: dict, first_unused: int, last: int) -> None:
    """Blank card/column slots that received no content.

    An unfilled card placeholder still renders its inherited card fill as an
    empty visible box, so clearing the text is not enough — the cloned
    placeholder shapes are removed from the slide entirely (the layout's
    originals are untouched). Handles both icon-card slots (header, card,
    subcopy, icon shape, icon) and text-column slots (header, card, subcopy).
    """
    for slot in range(first_unused, last + 1):
        for base in ("header", "card", "subcopy", "icon shape", "icon"):
            ph = _get_ph(slide, f"{base} {slot}", n2i)
            if ph is not None:
                ph._element.getparent().remove(ph._element)


# ── Layout handlers ────────────────────────────────────────────────────────

def handle_title_slide(prs, spec):
    layout = spec.get("layout", "Title slide_1 speaker")
    slide, n = add_slide(prs, layout)
    # Layout has: title (idx 0), subtitle (idx 1), name 1 (idx 22/23), company 1 (idx 24)
    fill_first(slide, ["title", "Title 3"], spec.get("title", ""), n)
    fill_first(slide, ["subtitle"], spec.get("subtitle", ""), n)
    fill_first(slide, ["name 1"], spec.get("presenter_name", ""), n)
    fill_first(slide, ["company 1"], spec.get("presenter_company", ""), n)
    return slide


def handle_section_header(prs, spec):
    slide, n = add_slide(prs, "Section Header")
    fill_first(slide, ["Title 1", "title"], spec.get("title", ""), n)
    return slide


def handle_eyebrow_content(prs, spec):
    layout = spec.get("layout", "Title+content+eyebrow_left")
    slide, n = add_slide(prs, layout)
    fill_first(slide, ["title"],   spec.get("title", ""),   n)
    fill_first(slide, ["eyebrow"], spec.get("eyebrow", ""), n)
    fill_content(slide, "content placeholder", spec.get("content", []), n)
    return slide


def handle_title_content(prs, spec):
    layout = spec.get("layout", "Title+content_left")
    slide, n = add_slide(prs, layout)
    fill_first(slide, ["title"], spec.get("title", ""), n)
    fill_content(slide, "content placeholder", spec.get("content", []), n)
    return slide


# Slot layouts in the template, by capacity. There is no icon-card layout
# smaller than 3 cards or larger than 6, and no text-column layout other
# than 2/3/4/6 columns. Both families are sized by CONTENT: the smallest
# layout that holds every item wins, so nothing is dropped while a larger
# layout exists and no empty slot renders when a smaller one fits.
ICON_CARD_LAYOUTS = {
    "3 icon cards+title": 3,
    "4 icon cards+title": 4,
    "icon cards+title":   6,
}
TEXT_COLUMN_LAYOUTS = {
    "2 text columns": 2,
    "3 text columns": 3,
    "4 text columns": 4,
    "6 text columns": 6,
}
# Any "N icon cards+title" / "N text columns" spelling routes to the slot
# handlers, whether or not N names a real template layout.
_ICON_CARD_RE    = re.compile(r"^(\d+ )?icon cards\+title$")
_TEXT_COLUMNS_RE = re.compile(r"^(\d+ )?text columns$")


def _fit_layout(n_items: int, layouts: dict) -> tuple:
    """(layout_name, capacity): the smallest layout holding every item, or
    the largest layout when nothing fits."""
    fits = [(cap, name) for name, cap in layouts.items() if cap >= n_items]
    cap, name = min(fits) if fits else max((cap, name) for name, cap in layouts.items())
    return name, cap


def _fill_slots(prs, spec, items, item_kind, layouts, family):
    """Shared body for the icon-card and text-column handlers."""
    requested = spec.get("layout")
    layout, capacity = _fit_layout(len(items), layouts)
    if requested and requested != layout:
        print(f"  ⚠ {family}: '{requested}' requested for {len(items)} "
              f"{item_kind}(s) — using '{layout}' (holds {capacity}) so the "
              f"content fits without empty slots.", file=sys.stderr)

    slide, n = add_slide(prs, layout)
    fill_first(slide, ["title"], spec.get("title", ""), n)
    for i, item in enumerate(items[:capacity], 1):
        fill_first(slide, [f"header {i}"],  item.get("header",  ""), n)
        fill_first(slide, [f"card {i}"],    item.get("body",    ""), n)
        fill_first(slide, [f"subcopy {i}"], item.get("subcopy", ""), n)

    if len(items) > capacity:
        print(f"  ⚠ {family}: {len(items)} {item_kind}(s) supplied but the "
              f"largest layout '{layout}' holds {capacity} — filled {capacity}, "
              f"DROPPED {len(items) - capacity}. Split the content across two "
              f"slides.", file=sys.stderr)
    elif len(items) < capacity:
        # The smallest layout that fits still has spare slots (the template
        # has no 1-, 2-, or 5-slot variants) — remove them rather than
        # render empty boxes.
        print(f"  ⚠ {family}: {len(items)} {item_kind}(s) on '{layout}' "
              f"({capacity} slots) — removed {capacity - len(items)} unused "
              f"slot(s).", file=sys.stderr)
        _blank_unused_card_slots(slide, n, len(items) + 1, capacity)
    return slide


def handle_icon_cards(prs, spec):
    return _fill_slots(prs, spec, spec.get("cards", []), "card",
                       ICON_CARD_LAYOUTS, "handle_icon_cards")


def handle_text_columns(prs, spec):
    return _fill_slots(prs, spec, spec.get("columns", []), "column",
                       TEXT_COLUMN_LAYOUTS, "handle_text_columns")


def handle_quote(prs, spec):
    slide, n = add_slide(prs, "Quote")
    fill_first(slide, ["title", "Title 1", "content placeholder"],
               spec.get("quote", ""), n)
    fill_first(slide, ["subtitle", "name 1", "company 1"],
               spec.get("attribution", ""), n)
    return slide


def handle_customer_story(prs, spec):
    layout = spec.get("layout", "Customer story")
    slide, n = add_slide(prs, layout)
    fill_first(slide, ["title", "Title 1"],                    spec.get("title", ""),   n)
    fill_content(slide, "content placeholder",                  spec.get("content", []), n)
    fill_first(slide, ["subtitle"],                             spec.get("subtitle", ""), n)
    return slide


def handle_chart(prs, spec):
    """Add a slide with a branded chart built from spec data.

    Spec keys:
        layout        — "Chart" (the dispatch key; not a template layout name)
        slide_layout  — real template layout to place the chart on
                        (default "Blank_graphic")
        title         — optional slide title (filled into the title placeholder)
        chart         — dict with:
            type          — XL_CHART_TYPE name string (default "BAR_CLUSTERED")
            categories    — list of category labels
            series        — list of {name, values} dicts
            left, top     — position in EMUs (default: 0.5in left / 1.5in top inset)
            width, height — size in EMUs (default: slide minus 1.0in / 2.0in insets)
    """
    # "Chart" is a dispatch key, not a template layout — resolve to a real one.
    layout = spec.get("slide_layout", "Blank_graphic")
    slide, n = add_slide(prs, layout)

    if spec.get("title"):
        fill_first(slide, ["title", "Title 1", "Title 3"], spec["title"], n)

    # Missing data is a failed slide, not a silent empty one: generate()
    # counts the ValueError and exits 1, so an incomplete deck never reports
    # success.
    chart_spec = spec.get("chart")
    if not chart_spec:
        raise ValueError("handle_chart: spec has no 'chart' object")

    # Resolve chart type — fall back to BAR_CLUSTERED on unknown strings
    type_name = chart_spec.get("type", "BAR_CLUSTERED").upper()
    chart_type = getattr(XL_CHART_TYPE, type_name, None)
    if chart_type is None:
        print(f"  ⚠ Unknown chart type '{type_name}' — falling back to BAR_CLUSTERED",
              file=sys.stderr)
        chart_type = XL_CHART_TYPE.BAR_CLUSTERED

    categories = chart_spec.get("categories", [])
    series_list = chart_spec.get("series", [])
    if not categories or not series_list:
        raise ValueError("handle_chart: 'chart.categories' and 'chart.series' "
                         "are both required")
    chart_data = ChartData()
    chart_data.categories = categories
    for s in series_list:
        chart_data.add_series(s.get("name", ""), s.get("values", []))

    # Default position/size: small insets on all four sides (inches → EMUs)
    slide_w, slide_h = _slide_dimensions(slide)
    default_left   = Inches(0.5)
    default_top    = Inches(1.5)
    default_width  = slide_w - Inches(1.0)
    default_height = slide_h - Inches(2.0)

    x  = Emu(chart_spec["left"])   if "left"   in chart_spec else default_left
    y  = Emu(chart_spec["top"])    if "top"    in chart_spec else default_top
    cx = Emu(chart_spec["width"])  if "width"  in chart_spec else default_width
    cy = Emu(chart_spec["height"]) if "height" in chart_spec else default_height

    graphic_frame = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)
    apply_brand_chart_colors(graphic_frame.chart)
    return slide


def handle_blank(prs, spec):
    layout = spec.get("layout", "Blank_graphic")
    slide, _ = add_slide(prs, layout)
    return slide


def handle_thank_you(prs, spec):
    slide, _ = add_slide(prs, "Thank you slide")
    return slide


def handle_generic(prs, spec):
    """Fallback: add the slide and fill any explicitly listed placeholders."""
    layout_name = spec.get("layout", "Title+content_left")
    slide, n = add_slide(prs, layout_name)
    for ph_name, text in spec.get("placeholders", {}).items():
        fill_content(slide, ph_name, text, n)
    return slide


# ── Dispatch ───────────────────────────────────────────────────────────────

HANDLERS = {
    "Title Slide":              handle_title_slide,
    "Title slide_1 speaker":    handle_title_slide,
    "Title slide_2 speakers":   handle_title_slide,
    "Title slide_3 speakers":   handle_title_slide,
    "Title slide_4 speakers":   handle_title_slide,
    "Section Header":                            handle_section_header,
    "Title+content+eyebrow_left":                handle_eyebrow_content,
    "Title+content+eyebrow_centered":            handle_eyebrow_content,
    "Title+content+eyebrow_middle aligned_left": handle_eyebrow_content,
    "Title+content+eyebrow_middle aligned_centered": handle_eyebrow_content,
    "Title+content_left":                        handle_title_content,
    "Title+content_centered":                    handle_title_content,
    "1_Title+content_left_2column":              handle_title_content,
    "3 icon cards+title":   handle_icon_cards,
    "4 icon cards+title":   handle_icon_cards,
    "icon cards+title":     handle_icon_cards,
    "2 text columns":       handle_text_columns,
    "3 text columns":       handle_text_columns,
    "4 text columns":       handle_text_columns,
    "6 text columns":       handle_text_columns,
    "Quote":                handle_quote,
    "Customer story":       handle_customer_story,
    "Customer story_stats": handle_customer_story,
    "Customer story_quote": handle_customer_story,
    "Blank_graphic":        handle_blank,
    "Blank_black":          handle_blank,
    "Chart":                handle_chart,
    "Thank you slide":      handle_thank_you,
}


def _route(layout: str):
    """Pick the handler for a spec's layout string, or None for the generic
    fallback. Exact names first, then the slot-layout families by pattern
    (so '5 text columns' reaches the content-sizing logic instead of failing
    as an unknown template layout), then prefix matches."""
    handler = HANDLERS.get(layout)
    if handler:
        return handler
    if _ICON_CARD_RE.match(layout):
        return handle_icon_cards
    if _TEXT_COLUMNS_RE.match(layout):
        return handle_text_columns
    for key, fn in HANDLERS.items():
        if layout.startswith(key):
            return fn
    return None


def dispatch(prs, spec: dict):
    layout = spec.get("layout", "")
    handler = _route(layout)
    if handler:
        slide = handler(prs, spec)
    else:
        print(f"  ⚠ No handler for '{layout}' — using generic fallback", file=sys.stderr)
        slide = handle_generic(prs, spec)

    # Apply wave background after the slide's content is placed
    if slide is not None and spec.get("wave_background"):
        wave = spec["wave_background"]
        opacity = float(spec.get("wave_overlay_opacity", 0.80))
        add_wave_background(slide, wave, opacity)

    return slide


# ── Generator ──────────────────────────────────────────────────────────────

def generate(spec_data: dict, output_path: Path) -> None:
    prs = load_template()
    clear_sample_slides(prs)

    slides_spec = spec_data.get("slides", [])
    if not slides_spec:
        print("Warning: spec contains no slides.", file=sys.stderr)

    failures = 0
    for i, slide_spec in enumerate(slides_spec):
        # Skip comment-only entries
        if set(slide_spec.keys()) <= {"_comment"}:
            continue
        layout = slide_spec.get("layout", "(none)")
        try:
            dispatch(prs, slide_spec)
            print(f"  ✓  [{i+1:2d}] {layout}")
        except Exception as e:
            print(f"  ✗  [{i+1:2d}] {layout} — {e}", file=sys.stderr)
            failures += 1

    if failures:
        print(f"\n⚠ {failures} slide(s) failed. Saving partial output.", file=sys.stderr)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"\nSaved → {output_path}")
    if failures:
        sys.exit(1)


# ── CLI ────────────────────────────────────────────────────────────────────

def main():
    if len(sys.argv) == 2 and sys.argv[1] == "--list-layouts":
        prs = load_template()
        print(f"{'Index':>5}  Layout name")
        print(f"{'─'*5}  {'─'*45}")
        for i, name in enumerate(list_layouts(prs)):
            print(f"{i:5d}  {name}")
        return

    if len(sys.argv) == 2 and sys.argv[1] == "--install-fonts":
        n = ensure_dtflow_fonts(verbose=True)
        if n == 0:
            print("No new fonts installed (already up to date).")
        return

    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    spec_path = Path(sys.argv[1])
    if not spec_path.exists():
        print(f"Spec file not found: {spec_path}", file=sys.stderr)
        sys.exit(1)

    output_path = Path(sys.argv[2]) if len(sys.argv) >= 3 \
        else spec_path.parent / f"deck-{date.today():%Y-%m-%d}.pptx"

    with open(spec_path, encoding="utf-8") as f:
        spec_data = json.load(f)

    check_dtflow_fonts()
    print(f"Template : {TEMPLATE.name}")
    print(f"Spec     : {spec_path}")
    print(f"Output   : {output_path}")
    print()
    generate(spec_data, output_path)


if __name__ == "__main__":
    main()
